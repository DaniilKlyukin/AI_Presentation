import os
from pptx import Presentation

class PptxToPpc:
    def __init__(self, pptx_path):
        self.prs = Presentation(pptx_path)

    def extract(self, output_txt_path=None):
        lines = []

        for idx, slide in enumerate(self.prs.slides):
            lines.append(f"# S:{slide.slide_id} (idx:{idx + 1})")

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes = slide.notes_slide.notes_text_frame.text.replace('\n', ' ')
                lines.append(f"> notes: {notes}")

            for shape in slide.shapes:
                if not (shape.has_text_frame or shape.has_table):
                    continue

                lines.append(f"\n## E:{shape.shape_id} ({shape.name})")

                if hasattr(shape, "left") and shape.left is not None:
                    l = int(shape.left / 12700)
                    t = int(shape.top / 12700)
                    w = int(shape.width / 12700)
                    h = int(shape.height / 12700)
                    lines.append(f"G:{l},{t},{w},{h}")

                if shape.has_text_frame and shape.text_frame.text.strip():
                    for p in shape.text_frame.paragraphs:
                        text = p.text.replace('\n', ' ').strip()
                        if not text: continue

                        style = f"lvl:{p.level}"
                        if p.runs and p.runs[0].font.size:
                            style += f",sz:{int(p.runs[0].font.size.pt)}"
                        if p.runs and p.runs[0].font.bold:
                            style += ",b:1"

                        lines.append(f"- [{style}] {text}")

                if shape.has_table:
                    for r_idx, row in enumerate(shape.table.rows):
                        row_cells = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text.replace('\n', ' ').strip()
                            row_cells.append(cell_text)

                        lines.append("| " + " | ".join(row_cells) + " |")
                        if r_idx == 0:
                            lines.append("|" + "|".join(["---"] * len(row.cells)) + "|")

            lines.append("\n" + "=" * 40 + "\n")

        result_text = '\n'.join(lines)

        if output_txt_path:
            with open(output_txt_path, 'w', encoding='utf-8') as f:
                f.write(result_text)

        return result_text