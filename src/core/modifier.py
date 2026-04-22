import os
import re
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from .errors import PptxSyntaxError, PptxLogicError


class CompactPptxModifier:
    def __init__(self, original_pptx):
        self.prs = Presentation(original_pptx)
        self.blank_layout = self._substitute_layout()
        self.DEFAULT_FONT = os.getenv("PPT_FONT", "Bookman Old Style")

        # Стандартные размеры слайда для проверки выхода за границы (в EMU)
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        self.warnings = []  # Сборник некритичных проблем

    def _substitute_layout(self):
        for layout in self.prs.slide_layouts:
            name = layout.name.lower()
            if "blank" in name or "пустой" in name: return layout
        return self.prs.slide_layouts[0] if len(self.prs.slide_layouts) > 0 else None

    def apply_from_file(self, ppc_txt_path, output_path):
        """Читает из файла и применяет изменения."""
        with open(ppc_txt_path, 'r', encoding='utf-8') as f:
            ppc_text = f.read()
        return self.apply_from_text(ppc_text, output_path)

    def apply_from_text(self, ppc_text, output_path):
        """Применяет изменения напрямую из текстовой переменной (для агента)."""
        self.warnings = []

        slide_pattern = re.compile(r'^#?\s*S:(\d+|NEW)(.*?)(?=^#?\s*S:|\Z)', re.MULTILINE | re.DOTALL)
        matches = list(slide_pattern.finditer(ppc_text))

        if not matches:
            raise PptxSyntaxError("В тексте не найдено ни одного блока слайда. Ожидается формат '# S:ID'.")

        target_order_ids = []

        for match in matches:
            s_id_str = match.group(1)
            block_content = match.group(2).strip()

            if s_id_str == 'NEW':
                slide = self.prs.slides.add_slide(self.blank_layout)
                target_order_ids.append(slide.slide_id)
                self._patch_slide_content(slide, block_content, "NEW")
            else:
                try:
                    s_id = int(s_id_str)
                except ValueError:
                    raise PptxSyntaxError(f"Некорректный ID слайда: {s_id_str}. Должен быть 'NEW' или число.")

                slide = self._find_slide_by_id(s_id)
                if slide:
                    target_order_ids.append(s_id)
                    self._patch_slide_content(slide, block_content, s_id)
                else:
                    self.warnings.append(f"Слайд с ID {s_id} не найден. Пропущен.")

        self._delete_missing_slides(target_order_ids)
        self._reorder_slides(target_order_ids)
        self.prs.save(output_path)

        return {
            "slides_processed": len(matches),
            "warnings": self.warnings
        }

    def _patch_slide_content(self, slide, block_text, slide_ref):
        notes_match = re.search(r'^>?\s*notes: (.*)', block_text, re.MULTILINE)
        if notes_match and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes_match.group(1).strip()

        element_pattern = re.compile(r'^#*\s*E:(\d+|NEW)[^\n]*(.*?)(?=^#*\s*E:|\Z)', re.MULTILINE | re.DOTALL)
        el_matches = list(element_pattern.finditer(block_text))

        received_el_ids = []
        for el_match in el_matches:
            el_id_str = el_match.group(1)
            el_content = el_match.group(2).strip().split('\n')

            if el_id_str == 'NEW':
                geom = self._extract_geometry(el_content, slide_ref, "NEW")
                new_shape = slide.shapes.add_textbox(
                    geom.get('l', Emu(914400)), geom.get('t', Emu(914400)),
                    geom.get('w', Emu(5000000)), geom.get('h', Emu(1000000))
                )
                received_el_ids.append(new_shape.shape_id)
                self._update_shape_data(new_shape, el_content, slide_ref, "NEW")
            else:
                try:
                    el_id = int(el_id_str)
                except ValueError:
                    raise PptxSyntaxError(f"Слайд {slide_ref}: Некорректный ID элемента: {el_id_str}.")

                shape = self._find_shape_by_id(slide, el_id)
                if shape:
                    received_el_ids.append(shape.shape_id)
                    self._update_shape_data(shape, el_content, slide_ref, el_id)
                else:
                    self.warnings.append(f"Слайд {slide_ref}: Элемент с ID {el_id} не найден. Пропущен.")

        for shape in list(slide.shapes):
            if (shape.has_text_frame or shape.has_table) and shape.shape_id not in received_el_ids:
                sp = shape._element
                sp.getparent().remove(sp)

    def _extract_geometry(self, lines, slide_ref, el_ref):
        for line in lines:
            line = line.strip()
            if line.startswith('G:'):
                try:
                    vals = list(map(int, line[2:].split(',')))
                    if len(vals) != 4:
                        raise ValueError()

                    l, t, w, h = [v * 12700 for v in vals]

                    # Проверка выхода за границы (Предупреждения для ИИ)
                    if l < 0 or t < 0 or (l + w) > self.slide_width or (t + h) > self.slide_height:
                        self.warnings.append(
                            f"Слайд {slide_ref}, Элемент {el_ref}: Координаты G:{line[2:]} выходят за границы слайда."
                        )

                    return {'l': Emu(l), 't': Emu(t), 'w': Emu(w), 'h': Emu(h)}
                except ValueError:
                    raise PptxSyntaxError(
                        f"Слайд {slide_ref}, Элемент {el_ref}: Синтаксическая ошибка в геометрии '{line}'. "
                        "Ожидается формат 'G:left,top,width,height' в целых числах."
                    )
        return {}

    def _update_shape_data(self, shape, data_lines, slide_ref, el_ref):
        geom = self._extract_geometry(data_lines, slide_ref, el_ref)
        if geom:
            shape.left, shape.top = geom['l'], geom['t']
            shape.width, shape.height = geom['w'], geom['h']

        paragraphs_data = []
        table_data = []

        for line in data_lines:
            line = line.strip()
            if not line: continue
            if line.startswith('G:'): continue
            if re.match(r'^\(.*\)$', line): continue

            if line.startswith('|'):
                if not re.match(r'^\|[-\s|]+\|$', line):
                    row_cells = [c.strip() for c in line.split('|')[1:-1]]
                    table_data.append(row_cells)
                continue

            base_style = {}
            block_style_match = re.search(r'\[((?:lvl|sz|b|c|i)[^\]]*)\]', line)
            if block_style_match:
                try:
                    base_style = dict(item.split(':') for item in block_style_match.group(1).split(',') if ':' in item)
                except ValueError:
                    raise PptxSyntaxError(
                        f"Слайд {slide_ref}, Элемент {el_ref}: Неверный синтаксис базового стиля '{block_style_match.group(0)}'.")

            clean_text = re.sub(r'\[(?:lvl|sz|b|c|i)[^\]]*\]', '', line)
            clean_text = re.sub(r'^-?\s*', '', clean_text).strip()

            if clean_text:
                paragraphs_data.append({'t': clean_text, 's': base_style})

        if shape.has_text_frame and paragraphs_data:
            self._apply_text(shape.text_frame, paragraphs_data, slide_ref, el_ref)
        if shape.has_table and table_data:
            self._apply_table(shape.table, table_data)

    def _parse_color(self, hex_str, slide_ref, el_ref):
        try:
            hc = hex_str.lstrip('#')
            if len(hc) != 6: raise ValueError()
            return RGBColor(int(hc[0:2], 16), int(hc[2:4], 16), int(hc[4:6], 16))
        except ValueError:
            self.warnings.append(
                f"Слайд {slide_ref}, Элемент {el_ref}: Неверный HEX цвет '{hex_str}'. Применен цвет по умолчанию.")
            return None

    def _apply_inline_formatting(self, paragraph, text, base_style, slide_ref, el_ref):
        paragraph.text = ""

        if 'lvl' in base_style:
            try:
                paragraph.level = min(int(base_style['lvl']), 8)
            except ValueError:
                raise PptxSyntaxError(f"Слайд {slide_ref}, Элемент {el_ref}: Уровень 'lvl' должен быть числом.")

        base_sz = Pt(int(base_style['sz'])) if 'sz' in base_style and base_style['sz'].isdigit() else None
        cur_b = base_style.get('b') == '1'
        cur_i = False
        cur_c = self._parse_color(base_style['c'], slide_ref, el_ref) if 'c' in base_style else None

        parts = re.split(r'(\[[bic]:[a-zA-Z0-9#]+\])', text)

        for part in parts:
            if not part: continue

            tag_match = re.match(r'^\[([bic]):([a-zA-Z0-9#]+)\]$', part)
            if tag_match:
                tag_type, tag_val = tag_match.group(1), tag_match.group(2)
                if tag_type == 'b':
                    cur_b = (tag_val == '1')
                elif tag_type == 'i':
                    cur_i = (tag_val == '1')
                elif tag_type == 'c':
                    cur_c = None if tag_val == '0' else self._parse_color(tag_val, slide_ref, el_ref)
                continue

            run = paragraph.add_run()
            run.text = part
            run.font.name = self.DEFAULT_FONT
            run.font.bold = cur_b
            run.font.italic = cur_i
            if base_sz: run.font.size = base_sz
            if cur_c: run.font.color.rgb = cur_c

    def _apply_text(self, tf, data, slide_ref, el_ref):
        tf.clear()
        for i, item in enumerate(data):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            self._apply_inline_formatting(p, item['t'], item['s'], slide_ref, el_ref)

    def _apply_table(self, table, table_data):
        # Логика таблицы оставлена без изменений
        for r_idx, row in enumerate(table_data):
            if r_idx < len(table.rows):
                for c_idx, text in enumerate(row):
                    if c_idx < len(table.columns):
                        cell = table.cell(r_idx, c_idx)
                        p = cell.text_frame.paragraphs[0] if len(
                            cell.text_frame.paragraphs) > 0 else cell.text_frame.add_paragraph()
                        self._apply_inline_formatting(p, text, {}, "Table", "Cell")

    def _reorder_slides(self, target_order):
        sldIdLst = self.prs.slides._sldIdLst
        id_to_xml = {slide.slide_id: sld_id for slide, sld_id in zip(self.prs.slides, sldIdLst)}
        for i, s_id in enumerate(target_order):
            if s_id in id_to_xml:
                xml_element = id_to_xml[s_id]
                sldIdLst.remove(xml_element)
                sldIdLst.insert(i, xml_element)

    def _delete_missing_slides(self, target_order):
        target_set = set(target_order)
        slides_to_remove = [i for i, slide in enumerate(self.prs.slides) if slide.slide_id not in target_set]
        for i in sorted(slides_to_remove, reverse=True):
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[i])

    def _find_slide_by_id(self, s_id):
        for slide in self.prs.slides:
            if slide.slide_id == s_id: return slide
        return None

    def _find_shape_by_id(self, slide, sh_id):
        for shape in slide.shapes:
            if shape.shape_id == sh_id: return shape
        return None