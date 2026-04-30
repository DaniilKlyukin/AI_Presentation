import os
import re
import uuid
from pathlib import Path
from dotenv import load_dotenv

import marko
from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml import parse_xml
from pygments import lexers
from pygments.lexers import get_lexer_by_name
import urllib.request
import ssl
from lxml import etree
import latex2mathml.converter

# Кастомные исключения
from .errors import PptxSyntaxError

env_path = Path(__file__).resolve().parent.parent.parent / '.env'
load_dotenv(dotenv_path=env_path)

CODE_THEME = {
    'Keyword': RGBColor(0, 0, 255),
    'Name': RGBColor(0, 0, 0),
    'Name.Function': RGBColor(121, 94, 38),
    'Name.Class': RGBColor(38, 127, 153),
    'String': RGBColor(163, 21, 21),
    'Comment': RGBColor(0, 128, 0),
    'Operator': RGBColor(0, 0, 0),
    'Number': RGBColor(9, 134, 88),
    'Type': RGBColor(38, 127, 153),
    'Other': RGBColor(0, 0, 0)
}


class PptxCreator:
    def __init__(self):
        self.prs = Presentation()
        self.font_name = os.getenv("PPT_FONT", "Bookman Old Style")
        self.title_size = int(os.getenv("PPT_TITLE_SIZE", "30"))
        self.body_size = int(os.getenv("PPT_BODY_SIZE", "22"))
        self.code_font = os.getenv("PPT_CODE_FONT", "Courier New")
        self.code_size = int(os.getenv("PPT_CODE_SIZE", "20"))
        self.line_spacing = float(os.getenv("PPT_LINE_SPACING", "1.0"))

        self.tittle_margin_cm = float(os.getenv("PPT_TITTLE_MARGIN_CM", "0.0"))
        self.body_margin_cm = float(os.getenv("PPT_BODY_MARGIN_CM", "0.0"))
        self.title_bg_color = os.getenv("PPT_TITLE_BG_COLOR", "")
        self.title_font_color = os.getenv("PPT_TITLE_FONT_COLOR", "0,0,0")
        self.title_height_cm = float(os.getenv("PPT_TITLE_HEIGHT_CM", "1.5"))

        self.slide_numbering = os.getenv("PPT_SLIDE_NUMBERING", "false").lower() == "true"
        self.footer_text = os.getenv("PPT_FOOTER_TEXT", "")
        self.footer_height_cm = float(os.getenv("PPT_FOOTER_HEIGHT_CM", "1.0"))
        self.formula_numbering = os.getenv("PPT_FORMULA_NUMBERING", "false").lower() == "true"
        self.bullet_spacing = float(os.getenv("PPT_BULLET_SPACING", "12.0"))

        self.footer_font_size = int(os.getenv("PPT_FOOTER_FONT_SIZE", "12"))
        self.numbering_font_size = int(os.getenv("PPT_NUMBERING_FONT_SIZE", "14"))
        self.numbering_width_cm = float(os.getenv("PPT_NUMBERING_WIDTH_CM", "2.0"))

        self.layout_title_idx = int(os.getenv("PPT_LAYOUT_TITLE_IDX", "0"))
        self.layout_content_idx = int(os.getenv("PPT_LAYOUT_CONTENT_IDX", "1"))
        self.content_bottom_buffer = float(os.getenv("PPT_CONTENT_BOTTOM_BUFFER_INCH", "1.2"))
        self.ts_title_top = float(os.getenv("PPT_TITLE_SLIDE_TITLE_TOP_INCH", "2.0"))
        self.ts_subtitle_top = float(os.getenv("PPT_TITLE_SLIDE_SUBTITLE_TOP_INCH", "4.0"))
        self.tf_padding = float(os.getenv("PPT_TEXT_FRAME_PADDING_CM", "0.13"))
        self.img_width_ratio = float(os.getenv("PPT_IMAGE_WIDTH_RATIO", "0.9"))
        self.table_row_h = float(os.getenv("PPT_TABLE_ROW_HEIGHT_CM", "1.0"))
        self.footer_border_color = os.getenv("PPT_FOOTER_BORDER_COLOR", "38,70,115")

        self.text_left_indent = float(os.getenv("PPT_TEXT_LEFT_INDENT_CM", "0.0"))
        self.text_first_line_indent = float(os.getenv("PPT_TEXT_FIRST_LINE_INDENT_CM", "0.0"))

        self._set_aspect_ratio(os.getenv("PPT_ASPECT_RATIO", "16:9"))
        self.title_layout = self.prs.slide_layouts[self.layout_title_idx]
        self.content_layout = self.prs.slide_layouts[self.layout_content_idx]

        self.formula_counter = 0
        self.warnings = []
        self._math_registry = {}
        self.md_parser = marko.Markdown(extensions=['gfm'])
        self._xslt_cache = None
        self.m_ns = "http://schemas.openxmlformats.org/officeDocument/2006/math"
        self._pending_math = {}  # ключ: placeholder_text, значение: omml_xml
        etree.register_namespace('m', self.m_ns)  # гарантирует префикс m: в XSLT-выводе

    def _parse_color(self, color_str, default_rgb=(0, 0, 0)):
        if not color_str: return RGBColor(*default_rgb)
        try:
            parts = [int(x.strip()) for x in color_str.split(',')]
            return RGBColor(parts[0], parts[1], parts[2])
        except:
            return RGBColor(*default_rgb)

    def _set_aspect_ratio(self, ratio_str):
        ratios = {"4:3": (9144000, 6858000), "16:9": (12192000, 6858000)}
        w, h = ratios.get(ratio_str, ratios["16:9"])
        self.prs.slide_width, self.prs.slide_height = w, h

    def _get_xslt(self):
        if self._xslt_cache: return self._xslt_cache
        path = os.path.join(os.path.dirname(__file__), "MML2OMML.XSL")

        if not os.path.exists(path):
            url = "https://raw.githubusercontent.com/plutext/docx4j/master/docx4j-openxml-objects/src/main/resources/org/docx4j/convert/out/common/xslt/MML2OMML.XSL"
            try:
                ctx = ssl.create_default_context()
                ctx.check_hostname = False
                ctx.verify_mode = ssl.CERT_NONE
                with urllib.request.urlopen(url, context=ctx, timeout=10) as r:
                    content = r.read()
                    with open(path, 'wb') as f: f.write(content)
            except Exception as e:
                self.warnings.append(f"Failed to download XSLT for math: {str(e)}")
                return None

        try:
            self._xslt_cache = etree.parse(path)
            return self._xslt_cache
        except Exception as e:
            self.warnings.append(f"Failed to parse XSLT: {str(e)}")
            return None

    def _latex_to_omml(self, latex_str, size_pt):
        """ Конвертирует LaTeX в валидный объект PowerPoint OMML """
        try:
            # 1. LaTeX -> MathML
            mathml = latex2mathml.converter.convert(latex_str)
            if 'xmlns' not in mathml:
                mathml = mathml.replace('<math', '<math xmlns="http://www.w3.org/1998/Math/MathML"', 1)

            # 2. XSLT Transform
            xslt_tree = self._get_xslt()
            if xslt_tree is None:
                return None

            transform = etree.XSLT(xslt_tree)
            tree = etree.fromstring(mathml.encode('utf-8'))
            omml_tree = transform(tree)

            m_ns = self.m_ns
            omml_nodes = omml_tree.xpath('//m:oMath', namespaces={'m': m_ns})
            if not omml_nodes:
                return None

            node = omml_nodes[0]

            # 3. Размер шрифта: 1 pt = 2 единицы в OMML
            sz_val = str(int(size_pt * 2))

            for rPr in node.xpath('.//m:rPr', namespaces={'m': m_ns}):
                # Удаляем старые теги размера
                for old_sz in rPr.xpath('./m:sz', namespaces={'m': m_ns}):
                    rPr.remove(old_sz)
                # Добавляем правильный размер
                sz = etree.SubElement(rPr, "{%s}sz" % m_ns)
                sz.set("{%s}val" % m_ns, sz_val)
                # Шрифт Cambria Math (стандарт Office)
                rf = etree.SubElement(rPr, "{%s}rFonts" % m_ns)
                rf.set("{%s}ascii" % m_ns, "Cambria Math")
                rf.set("{%s}hAnsi" % m_ns, "Cambria Math")

            # 4. Сериализация с гарантированным префиксом m:
            xml_str = etree.tostring(node, encoding='unicode')
            # На всякий случай добавляем xmlns:m, если его нет (обычно уже есть)
            if 'xmlns:m' not in xml_str.split('>')[0]:
                xml_str = xml_str.replace('<m:oMath', f'<m:oMath xmlns:m="{m_ns}"', 1)

            return parse_xml(xml_str)
        except Exception as e:
            print(f"[ERROR] Ошибка конвертации: {e}")
            return None

    def _process_math_blocks(self, text):
        """ Заменяет формулы на маркеры и выводит инфо в консоль """
        self._math_registry = {}

        def repl(m):
            original = m.group(0)
            is_block = original.startswith('$$')
            latex = original.strip('$').strip()
            marker_id = f"MATHM{uuid.uuid4().hex}X"
            self._math_registry[marker_id] = (latex, is_block, original)
            return marker_id

        # Сначала блочные, потом строчные
        text = re.sub(r'\$\$.*?\$\$', repl, text, flags=re.DOTALL)
        text = re.sub(r'(?<!\$)\$(?!\$).*?\$', repl, text)

        print(f"[DEBUG] Найдено формул в блоке: {len(self._math_registry)}")
        for mid, data in self._math_registry.items():
            print(f"  - Маркер: {mid} | LaTeX: {data[0][:30]}...")

        return text

    def _apply_paragraph_style(self, paragraph, is_code=False, align=None, is_list_item=False):
        if align: paragraph.alignment = align
        paragraph.line_spacing = self.line_spacing
        if is_list_item:
            paragraph.space_after = Pt(self.bullet_spacing)
        else:
            paragraph.level = 0
            self._remove_bullet_xml(paragraph)

    def _remove_bullet_xml(self, paragraph):
        pPr = paragraph._p.get_or_add_pPr()
        for tag in ['buNone', 'buChar', 'buAutoNum']:
            for e in pPr.findall(f'.//{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}'):
                pPr.remove(e)
        pPr.insert(0, OxmlElement('a:buNone'))
        l_margin = int(Cm(self.text_left_indent).emu)
        indent = int(Cm(self.text_first_line_indent - self.text_left_indent).emu)
        pPr.set('marL', str(l_margin))
        pPr.set('indent', str(indent))

    def _setup_text_frame(self, shape, align=None, is_title=False):
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT if is_title else MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE if is_title else MSO_ANCHOR.TOP
        margin = Cm(self.tittle_margin_cm) if is_title else Cm(self.body_margin_cm)
        tf.margin_left = tf.margin_right = margin
        tf.margin_top = tf.margin_bottom = Cm(self.tf_padding)
        tf.clear()
        return tf

    def _position_shape(self, shape, top_inch, height_inch=None):
        shape.width = self.prs.slide_width - Cm(self.tittle_margin_cm * 2)
        if height_inch: shape.height = Inches(height_inch)
        shape.left = Cm(self.tittle_margin_cm)
        shape.top = Inches(top_inch)
        return shape

    def _get_or_add_paragraph(self, text_frame):
        if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
            return text_frame.paragraphs[0]
        return text_frame.add_paragraph()

    def _get_current_tf_height(self, text_frame):
        total_pt = 0
        for p in text_frame.paragraphs:
            max_size = self.body_size
            for run in p.runs:
                if run.font.size: max_size = max(max_size, run.font.size.pt)
            lines = (len(p.text) // max(40, int(1400 / max_size))) + 1
            total_pt += lines * max_size * self.line_spacing * 1.2
        return total_pt / 72.0

    def _insert_image_shape(self, slide, text_frame, img_path, level=0, is_block=True):
        if not os.path.exists(img_path): return
        current_h = self._get_current_tf_height(text_frame)
        top = text_frame._parent.top.inches + current_h + 0.1
        pic = slide.shapes.add_picture(img_path, Inches(0), Inches(top))
        max_w = (self.prs.slide_width - Cm(self.tittle_margin_cm * 2)) * self.img_width_ratio
        if pic.width > max_w:
            pic.height = int(pic.height * (max_w / pic.width))
            pic.width = int(max_w)
        pic.left = int((self.prs.slide_width - pic.width) / 2) if is_block else Cm(0.5)
        if is_block:
            p = text_frame.add_paragraph()
            p.font.size = Pt(1)
            p.space_before = Pt(pic.height.pt + 10)

    def _add_table(self, slide, text_frame, node):
        rows_nodes = node.children
        rows, cols = len(rows_nodes), max((len(r.children) for r in rows_nodes), default=0)
        if rows == 0 or cols == 0: return
        top = text_frame._parent.top.inches + self._get_current_tf_height(text_frame) + 0.2
        width = int(self.prs.slide_width * self.img_width_ratio)
        tbl_shape = slide.shapes.add_table(rows, cols, int((self.prs.slide_width - width) / 2), Inches(top), width,
                                           Cm(self.table_row_h * rows))
        for r_idx, row in enumerate(rows_nodes):
            for c_idx, cell in enumerate(row.children):
                if c_idx < cols:
                    cell_tf = tbl_shape.table.cell(r_idx, c_idx).text_frame
                    cell_tf.clear()
                    p = cell_tf.paragraphs[0]
                    p.alignment = PP_ALIGN.CENTER
                    self._fill_run(p, cell)
                    for run in p.runs:
                        run.font.name, run.font.size = self.font_name, Pt(self.body_size - 4)
                        if r_idx == 0: run.font.bold = True
        p = text_frame.add_paragraph()
        p.font.size = Pt(1)
        p.space_before = Pt(tbl_shape.height.pt + 15)

    def _add_highlighted_code(self, text_frame, code_text, lang, slide, title_node, slide_idx):
        # Восстанавливаем формулы в коде перед подсветкой
        code_text = self._restore_math_in_text(code_text)

        lexer = get_lexer_by_name(lang) if lang else lexers.get_lexer_by_name('text')
        for line in code_text.replace('\t', '    ').splitlines():
            footer_h = Cm(self.footer_height_cm).inches if (self.footer_text or self.slide_numbering) else 0
            if self._get_current_tf_height(text_frame) > (
                    self.prs.slide_height.inches - Cm(self.title_height_cm + 2).inches - footer_h):
                slide_idx += 1
                slide, text_frame = self._init_content_slide(title_node, slide_idx, is_continuation=True)
            p = self._get_or_add_paragraph(text_frame)
            self._apply_paragraph_style(p, is_code=True, align=PP_ALIGN.LEFT)

            for ttype, value in lexer.get_tokens(line):
                if not value.strip('\r\n'): continue
                run = p.add_run()
                run.text = value
                run.font.name, run.font.size = self.code_font, Pt(self.code_size)
                key = str(ttype).split('.')[-1]
                run.font.color.rgb = CODE_THEME.get(key, CODE_THEME['Other'])
        return text_frame, slide, slide_idx


    def _restore_math_in_text(self, text):
        """ Вспомогательный метод для восстановления исходного текста (для кода) """
        for mid, data in self._math_registry.items():
            text = text.replace(mid, data[2])
        return text


    def _add_node_to_frame(self, text_frame, node, slide=None, level=0, default_align=None, is_list_item=False,
                           is_quote=False, title_node=None, slide_idx=0):
        ntype = node.__class__.__name__

        def check_overflow(tf, s, idx):
            max_h = self.prs.slide_height.inches - Cm(self.title_height_cm + 2).inches - self.content_bottom_buffer
            if self._get_current_tf_height(tf) > max_h:
                new_s, new_tf = self._init_content_slide(title_node, idx + 1, True)
                return new_tf, new_s, idx + 1
            return tf, s, idx

        if ntype == 'Paragraph':
            text_frame, slide, slide_idx = check_overflow(text_frame, slide, slide_idx)
            if len(node.children) == 1 and node.children[0].__class__.__name__ == 'Image':
                self._insert_image_shape(slide, text_frame, node.children[0].dest, level, True)
            else:
                p = self._get_or_add_paragraph(text_frame)
                p.level = min(level, 8)
                self._apply_paragraph_style(p, align=default_align, is_list_item=is_list_item)
                for child in node.children: self._fill_run(p, child, italic=is_quote)
        elif ntype == 'Quote':
            for child in node.children:
                text_frame, slide, slide_idx = self._add_node_to_frame(text_frame, child, slide, level,
                                                                       title_node=title_node, slide_idx=slide_idx,
                                                                       is_quote=True)
        elif ntype == 'List':
            for item in node.children:
                for sub in item.children:
                    text_frame, slide, slide_idx = self._add_node_to_frame(text_frame, sub, slide, level + (
                        1 if sub.__class__.__name__ == 'List' else 0), default_align, True, is_quote, title_node,
                                                                           slide_idx)
        elif ntype in ['FencedCode', 'CodeBlock']:
            content = node.children[0].children if hasattr(node.children[0], 'children') else node.children[0]
            text_frame, slide, slide_idx = self._add_highlighted_code(text_frame, str(content),
                                                                      getattr(node, 'lang', 'text'), slide, title_node,
                                                                      slide_idx)
        elif ntype == 'Table':
            text_frame, slide, slide_idx = check_overflow(text_frame, slide, slide_idx)
            self._add_table(slide, text_frame, node)
        return text_frame, slide, slide_idx

    def _fill_run(self, paragraph, node, is_title=False, bold=False, italic=False):
        if node is None: return
        ntype = node.__class__.__name__

        cur_bold = bold or ntype in ['Strong', 'StrongEmphasis']
        cur_italic = italic or ntype in ['Emphasis', 'Italic']
        is_code = (ntype == 'CodeSpan')

        if isinstance(node, str):
            self._create_run(paragraph, node, is_title, cur_bold, cur_italic, is_code)
        elif ntype == 'RawText':
            self._create_run(paragraph, node.children, is_title, cur_bold, cur_italic, is_code)
        elif ntype == 'Image':
            alt = getattr(node, 'title', '') or (node.children[0].children if node.children else '')
            self._create_run(paragraph, str(alt), is_title, cur_bold, cur_italic)
        elif hasattr(node, 'children'):
            if isinstance(node.children, str):
                self._create_run(paragraph, node.children, is_title, cur_bold, cur_italic, is_code)
            else:
                for child in node.children:
                    self._fill_run(paragraph, child, is_title, cur_bold, cur_italic)
        elif ntype in ['LineBreak', 'SoftLineBreak']:
            paragraph.add_run().text = " "

    def _create_run(self, paragraph, text, is_title, bold=False, italic=False, is_code=False):
        if not text:
            return

        parts = re.split(r'(MATHM[a-f0-9]{32}X)', str(text))
        for part in parts:
            if part in self._math_registry:
                latex, is_block, _ = self._math_registry[part]

                placeholder_uuid = uuid.uuid4().hex
                placeholder_text = f"«{placeholder_uuid}»"

                omml = self._latex_to_omml(latex, self.title_size if is_title else self.body_size)
                if omml is not None:
                    # ВАЖНО: Сохраняем только сам omml, без обертки в a:r
                    self._pending_math[placeholder_text] = etree.tostring(omml, encoding='unicode')

                run = paragraph.add_run()
                run.text = placeholder_text
                self._apply_run_style(run, is_title, bold, italic, False)
            elif part:
                run = paragraph.add_run()
                run.text = part
                self._apply_run_style(run, is_title, bold, italic, is_code)

    def _apply_run_style(self, run, is_title, bold, italic, is_code):
        run.font.name = self.code_font if is_code else self.font_name
        size = self.code_size if is_code else (self.title_size if is_title else self.body_size)
        run.font.size = Pt(size)
        if is_title:
            run.font.color.rgb = self._parse_color(self.title_font_color)
        run.font.bold = (bold or is_title)
        run.font.italic = italic

    def _add_footer_and_numbering(self, slide, slide_idx):
        sw, sh = self.prs.slide_width, self.prs.slide_height
        margin, fh = Cm(self.tittle_margin_cm), Cm(self.footer_height_cm)
        by = sh - fh
        bc = self._parse_color(self.footer_border_color)
        nw = Cm(self.numbering_width_cm) if self.slide_numbering else 0

        if self.footer_text:
            fw = sw - margin * 2 - nw
            f_s = slide.shapes.add_textbox(margin, by, fw, fh)
            f_s.line.color.rgb, f_s.line.width = bc, Pt(1)
            p = f_s.text_frame.paragraphs[0]
            p.text, p.alignment = self.footer_text, PP_ALIGN.CENTER
            p.font.size, p.font.name = Pt(self.footer_font_size), self.font_name

        if self.slide_numbering:
            nl = margin + (sw - margin * 2 - nw) if self.footer_text else sw - margin - nw
            n_s = slide.shapes.add_textbox(nl, by, nw, fh)
            n_s.name = "SlideNumberBox"
            n_s.line.color.rgb, n_s.line.width = bc, Pt(1)
            p = n_s.text_frame.paragraphs[0]
            p.text, p.alignment = str(slide_idx), PP_ALIGN.CENTER
            p.font.size, p.font.bold, p.font.name = Pt(self.numbering_font_size), True, self.font_name

    def _inject_math_placeholders(self, pptx_path):
        import zipfile
        import shutil
        from lxml import etree

        tmp_path = pptx_path + ".tmp"
        zin = zipfile.ZipFile(pptx_path, 'r')
        zout = zipfile.ZipFile(tmp_path, 'w', zipfile.ZIP_DEFLATED)

        math_map = self._pending_math
        # Регистрируем пространства имен для корректной записи
        nsmap = {
            'a': "http://schemas.openxmlformats.org/drawingml/2006/main",
            'm': "http://schemas.openxmlformats.org/officeDocument/2006/math",
            'p': "http://schemas.openxmlformats.org/presentationml/2006/main"
        }

        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                root = etree.fromstring(data)

                # Ищем все текстовые прогоны
                for run in root.findall('.//a:r', namespaces=nsmap):
                    t_node = run.find('a:t', namespaces=nsmap)
                    if t_node is not None and t_node.text in math_map:
                        math_xml = math_map[t_node.text]
                        # Создаем элемент формулы из сохраненной строки
                        new_node = etree.fromstring(math_xml)

                        # Заменяем старый <a:r> на <m:oMath> в родителе (<a:p>)
                        parent = run.getparent()
                        if parent is not None:
                            parent.replace(run, new_node)

                # Сериализуем обратно
                content = etree.tostring(root, encoding='utf-8', xml_declaration=True)
                zout.writestr(item, content)
            else:
                zout.writestr(item, data)

        zin.close()
        zout.close()
        shutil.move(tmp_path, pptx_path)

    def _init_content_slide(self, title_node, slide_idx, is_continuation=False):
        slide = self.prs.slides.add_slide(self.content_layout)
        if slide.shapes.title and title_node:
            shape = self._position_shape(slide.shapes.title, Cm(self.tittle_margin_cm).inches,
                                         Cm(self.title_height_cm).inches)
            if self.title_bg_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._parse_color(self.title_bg_color)
            tf = self._setup_text_frame(shape, PP_ALIGN.CENTER, True)
            self._fill_run(tf.paragraphs[0], title_node, True)
        self._add_footer_and_numbering(slide, slide_idx)
        top = Cm(self.tittle_margin_cm + self.title_height_cm).inches
        h = self.prs.slide_height.inches - top - Cm(
            self.footer_height_cm if self.footer_text or self.slide_numbering else self.tittle_margin_cm).inches
        shape = self._position_shape(slide.placeholders[1], top, h)
        return slide, self._setup_text_frame(shape, PP_ALIGN.LEFT)

    def _create_title_slide(self, doc, slide_num):
        slide = self.prs.slides.add_slide(self.title_layout)
        h_node = next((n for n in doc.children if n.__class__.__name__ == 'Heading'), None)
        if slide.shapes.title and h_node:
            shape = self._position_shape(slide.shapes.title, self.ts_title_top, Cm(self.title_height_cm).inches)
            if self.title_bg_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._parse_color(self.title_bg_color)
            self._fill_run(self._setup_text_frame(shape, PP_ALIGN.CENTER, True).paragraphs[0], h_node, True)
        if len(slide.placeholders) > 1:
            shape = self._position_shape(slide.placeholders[1], self.ts_subtitle_top, 2.5)
            tf = self._setup_text_frame(shape, PP_ALIGN.CENTER)
            for node in [n for n in doc.children if n != h_node]:
                self._add_node_to_frame(tf, node, slide, default_align=PP_ALIGN.CENTER)
        self._add_footer_and_numbering(slide, slide_num)

    def create_from_text(self, md_text, output_path):
        self._pending_math.clear()
        self.warnings, self.formula_counter = [], 0
        md_text = self._process_math_blocks(md_text)  # Защищаем формулы перед marko
        blocks = [b for b in re.split(r'\n\s*---\s*\n', md_text.strip()) if b.strip()]
        if not blocks: raise PptxSyntaxError("MD текст пуст.")

        slide_num = 1
        for idx, block in enumerate(blocks):
            doc = self.md_parser.parse(block)
            if idx == 0:
                self._create_title_slide(doc, slide_num)
                slide_num += 1
            else:
                slide_num = self._create_content_slide(doc, slide_num) + 1

        total = len(self.prs.slides)
        if self.slide_numbering:
            for s in self.prs.slides:
                for sh in s.shapes:
                    if sh.name == "SlideNumberBox":
                        p = sh.text_frame.paragraphs[0]
                        p.text = f"{p.text}/{total}"
                        p.alignment, p.font.bold = PP_ALIGN.CENTER, True

        self.prs.save(output_path)

        self._inject_math_placeholders(output_path)

        return {"slides_created": total, "warnings": self.warnings}

    def _create_content_slide(self, doc, start_idx):
        title = next((n for n in doc.children if n.__class__.__name__ == 'Heading'), None)
        slide, tf = self._init_content_slide(title, start_idx)
        idx = start_idx
        for node in [n for n in doc.children if n != title]:
            tf, slide, idx = self._add_node_to_frame(tf, node, slide, title_node=title, slide_idx=idx)
        return idx

    def create_from_file(self, md_path, output_path):
        with open(md_path, 'r', encoding='utf-8') as f:
            return self.create_from_text(f.read(), output_path)