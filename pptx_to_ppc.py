# pptx_to_ppc.py
import os
from pptx import Presentation

class PptxToPpc:
    def __init__(self, pptx_path):
        self.prs = Presentation(pptx_path)

    def extract(self, output_txt_path):
        lines = []

        for idx, slide in enumerate(self.prs.slides):
            lines.append(f"# S:{slide.slide_id} (idx:{idx + 1})")

            # Извлекаем заметки
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes = slide.notes_slide.notes_text_frame.text.replace('\n', ' ')
                lines.append(f"> notes: {notes}")

            # Извлекаем элементы
            for shape in slide.shapes:
                # Фильтруем элементы без контента (линии, пустые группы)
                if not (shape.has_text_frame or shape.has_table):
                    continue

                lines.append(f"\n## E:{shape.shape_id} ({shape.name})")

                # Извлекаем геометрию (переводим из EMU в PT)
                if hasattr(shape, "left") and shape.left is not None:
                    l = int(shape.left / 12700)
                    t = int(shape.top / 12700)
                    w = int(shape.width / 12700)
                    h = int(shape.height / 12700)
                    lines.append(f"G:{l},{t},{w},{h}")

                # Извлекаем текст
                if shape.has_text_frame and shape.text_frame.text.strip():
                    for p in shape.text_frame.paragraphs:
                        text = p.text.replace('\n', ' ').strip()
                        if not text: continue

                        # Собираем базовые стили из первого прогона (run)
                        style = f"lvl:{p.level}"
                        if p.runs and p.runs[0].font.size:
                            style += f",sz:{int(p.runs[0].font.size.pt)}"
                        if p.runs and p.runs[0].font.bold:
                            style += ",b:1"

                        lines.append(f"- [{style}] {text}")

                # Извлекаем таблицы (переводим в Markdown)
                if shape.has_table:
                    for r_idx, row in enumerate(shape.table.rows):
                        row_cells = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text.replace('\n', ' ').strip()
                            row_cells.append(cell_text)

                        lines.append("| " + " | ".join(row_cells) + " |")
                        if r_idx == 0:
                            lines.append("|" + "|".join(["---"] * len(row.cells)) + "|")

            lines.append("\n" + "=" * 40 + "\n")

        with open(output_txt_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        print(f"✅ Данные успешно извлечены в: {output_txt_path}")

if __name__ == "__main__":
    pptx_file = input("Введите путь к исходному PPTX: ").strip('" ')
    if os.path.exists(pptx_file):
        out_file = pptx_file.replace('.pptx', '.txt')
        extractor = PptxToPpc(pptx_file)
        extractor.extract(out_file)
    else:
        print("❌ Файл не найден.")