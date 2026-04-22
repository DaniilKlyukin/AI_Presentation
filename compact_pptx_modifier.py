import os
import re
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pathlib import Path


class CompactPptxModifier:
    def __init__(self, original_pptx):
        self.prs = Presentation(original_pptx)
        self.blank_layout = self._substitute_layout()
        self.DEFAULT_FONT = os.getenv("PPT_FONT", "Bookman Old Style")

    def _substitute_layout(self):
        for layout in self.prs.slide_layouts:
            name = layout.name.lower()
            if "blank" in name or "пустой" in name:
                return layout
        return self.prs.slide_layouts[0] if len(self.prs.slide_layouts) > 0 else None

    def apply_ppc(self, ppc_txt_path, output_path):
        with open(ppc_txt_path, 'r', encoding='utf-8') as f:
            ppc_text = f.read()

        slide_pattern = re.compile(r'^#?\s*S:(\d+|NEW)(.*?)(?=^#?\s*S:|\Z)', re.MULTILINE | re.DOTALL)
        matches = list(slide_pattern.finditer(ppc_text))

        if not matches:
            print("❌ ОШИБКА: В текстовом файле не найдено ни одного блока слайда (S:ID).")
            return

        print(f"🔍 Найдено слайдов в файле: {len(matches)}")
        target_order_ids = []

        for match in matches:
            s_id_str = match.group(1)
            block_content = match.group(2).strip()

            if s_id_str == 'NEW':
                slide = self.prs.slides.add_slide(self.blank_layout)
                target_order_ids.append(slide.slide_id)
                self._patch_slide_content(slide, block_content)
            else:
                s_id = int(s_id_str)
                slide = self._find_slide_by_id(s_id)
                if slide:
                    target_order_ids.append(s_id)
                    self._patch_slide_content(slide, block_content)

        self._delete_missing_slides(target_order_ids)
        self._reorder_slides(target_order_ids)
        self.prs.save(output_path)
        print(f"✅ Готово! Результат сохранен в: {output_path}")

    def _patch_slide_content(self, slide, block_text):
        notes_match = re.search(r'^>?\s*notes: (.*)', block_text, re.MULTILINE)
        if notes_match and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes_match.group(1).strip()

        # ИСПРАВЛЕНИЕ 1: [^\n]* "съедает" техническое имя (TextBox 10) до конца строки,
        # чтобы оно не попало в контент (group 2).
        element_pattern = re.compile(r'^#*\s*E:(\d+|NEW)[^\n]*(.*?)(?=^#*\s*E:|\Z)', re.MULTILINE | re.DOTALL)
        el_matches = list(element_pattern.finditer(block_text))

        received_el_ids = []
        for el_match in el_matches:
            el_id_str = el_match.group(1)
            el_content = el_match.group(2).strip().split('\n')

            if el_id_str == 'NEW':
                geom = self._extract_geometry(el_content)
                new_shape = slide.shapes.add_textbox(
                    geom.get('l', Emu(914400)), geom.get('t', Emu(914400)),
                    geom.get('w', Emu(5000000)), geom.get('h', Emu(1000000))
                )
                received_el_ids.append(new_shape.shape_id)
                self._update_shape_data(new_shape, el_content)
            else:
                el_id = int(el_id_str)
                shape = self._find_shape_by_id(slide, el_id)
                if shape:
                    received_el_ids.append(shape.shape_id)
                    self._update_shape_data(shape, el_content)

        for shape in list(slide.shapes):
            if (shape.has_text_frame or shape.has_table) and shape.shape_id not in received_el_ids:
                sp = shape._element
                sp.getparent().remove(sp)

    def _extract_geometry(self, lines):
        for line in lines:
            line = line.strip()
            if line.startswith('G:'):
                try:
                    vals = list(map(int, line[2:].split(',')))
                    return {'l': Emu(vals[0] * 12700), 't': Emu(vals[1] * 12700), 'w': Emu(vals[2] * 12700),
                            'h': Emu(vals[3] * 12700)}
                except:
                    pass
        return {}

    def _update_shape_data(self, shape, data_lines):
        geom = self._extract_geometry(data_lines)
        if geom:
            shape.left, shape.top = geom['l'], geom['t']
            shape.width, shape.height = geom['w'], geom['h']

        paragraphs_data = []
        table_data = []

        for line in data_lines:
            line = line.strip()
            if not line: continue

            # ИСПРАВЛЕНИЕ 2: Жестко игнорируем служебные строки ИИ, если они просочились
            if line.startswith('G:'):
                continue  # Игнорируем координаты
            if re.match(r'^\(.*\)$', line):
                continue  # Игнорируем названия вроде (Rectangle 28)

            # Обработка таблиц
            if line.startswith('|'):
                if not re.match(r'^\|[-\s|]+\|$', line):
                    row_cells = [c.strip() for c in line.split('|')[1:-1]]
                    table_data.append(row_cells)
                continue

            # Обработка текстовых строк (Ищем базовые стили: lvl, sz)
            base_style = {}
            block_style_match = re.search(r'\[((?:lvl|sz)[^\]]*)\]', line)
            if block_style_match:
                base_style = dict(item.split(':') for item in block_style_match.group(1).split(',') if ':' in item)

            # Удаляем базовые теги и начальные дефисы из текста
            clean_text = re.sub(r'\[(?:lvl|sz)[^\]]*\]', '', line)
            clean_text = re.sub(r'^-?\s*', '', clean_text).strip()

            if clean_text:
                paragraphs_data.append({'t': clean_text, 's': base_style})

        if shape.has_text_frame and paragraphs_data:
            self._apply_text(shape.text_frame, paragraphs_data)
        if shape.has_table and table_data:
            self._apply_table(shape.table, table_data)

    def _parse_color(self, hex_str):
        try:
            hc = hex_str.lstrip('#')
            return RGBColor(int(hc[0:2], 16), int(hc[2:4], 16), int(hc[4:6], 16))
        except:
            return None

    def _apply_inline_formatting(self, paragraph, text, base_style):
        paragraph.text = ""

        if 'lvl' in base_style: paragraph.level = min(int(base_style['lvl']), 8)
        base_sz = Pt(int(base_style['sz'])) if 'sz' in base_style else None

        cur_b = base_style.get('b') == '1'
        cur_i = False
        cur_c = self._parse_color(base_style['c']) if 'c' in base_style else None

        parts = re.split(r'(\[[bic]:[a-zA-Z0-9#]+\])', text)

        for part in parts:
            if not part: continue

            tag_match = re.match(r'^\[([bic]):([a-zA-Z0-9#]+)\]$', part)
            if tag_match:
                tag_type, tag_val = tag_match.group(1), tag_match.group(2)
                if tag_type == 'b':
                    cur_b = (tag_val == '1')
                elif tag_type == 'i':
                    cur_i = (tag_val == '1')
                elif tag_type == 'c':
                    cur_c = None if tag_val == '0' else self._parse_color(tag_val)
                continue

            run = paragraph.add_run()
            run.text = part
            run.font.name = self.DEFAULT_FONT
            run.font.bold = cur_b
            run.font.italic = cur_i
            if base_sz: run.font.size = base_sz
            if cur_c: run.font.color.rgb = cur_c

    def _apply_text(self, tf, data):
        tf.clear()
        for i, item in enumerate(data):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            self._apply_inline_formatting(p, item['t'], item['s'])

    def _apply_table(self, table, table_data):
        for r_idx, row in enumerate(table_data):
            if r_idx < len(table.rows):
                for c_idx, text in enumerate(row):
                    if c_idx < len(table.columns):
                        cell = table.cell(r_idx, c_idx)
                        p = cell.text_frame.paragraphs[0] if len(
                            cell.text_frame.paragraphs) > 0 else cell.text_frame.add_paragraph()
                        self._apply_inline_formatting(p, text, {})

    def _reorder_slides(self, target_order):
        sldIdLst = self.prs.slides._sldIdLst
        id_to_xml = {slide.slide_id: sld_id for slide, sld_id in zip(self.prs.slides, sldIdLst)}
        for i, s_id in enumerate(target_order):
            if s_id in id_to_xml:
                xml_element = id_to_xml[s_id]
                sldIdLst.remove(xml_element)
                sldIdLst.insert(i, xml_element)

    def _delete_missing_slides(self, target_order):
        target_set = set(target_order)
        slides_to_remove = [i for i, slide in enumerate(self.prs.slides) if slide.slide_id not in target_set]
        for i in sorted(slides_to_remove, reverse=True):
            self.prs.slides._sldIdLst.remove(self.prs.slides._sldIdLst[i])

    def _find_slide_by_id(self, s_id):
        for slide in self.prs.slides:
            if slide.slide_id == s_id: return slide
        return None

    def _find_shape_by_id(self, slide, sh_id):
        for shape in slide.shapes:
            if shape.shape_id == sh_id: return shape
        return None


if __name__ == "__main__":
    pptx_in = input("Путь к PPTX: ").strip('" ')
    path_obj = Path(pptx_in)

    ppc_path_input = input("Путь к TXT (Enter если файл рядом): ").strip('" ')
    ppc_path = Path(ppc_path_input) if ppc_path_input else path_obj.with_suffix('.txt')

    out_path = path_obj.parent / f"modified_{path_obj.name}"

    if path_obj.exists() and ppc_path.exists():
        modifier = CompactPptxModifier(str(path_obj))
        modifier.apply_ppc(str(ppc_path), str(out_path))
    else:
        print(f"❌ Файлы не найдены.")